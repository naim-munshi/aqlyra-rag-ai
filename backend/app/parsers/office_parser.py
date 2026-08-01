from datetime import date, datetime, time
from pathlib import Path
from typing import Any

from docx import Document as DocxDocument
from docx.table import Table as DocxTable
from docx.text.paragraph import Paragraph
from openpyxl import load_workbook
from pptx import Presentation

from app.parsers.common import (
    calculate_quality_score,
    create_parsed_unit,
    normalize_text,
)
from app.parsers.types import (
    DocumentParsingError,
    ParseResult,
    ParsedUnit,
)


MAX_SPREADSHEET_ROWS = 50_000
MAX_SPREADSHEET_COLUMNS = 200


def _table_rows_to_text(
    rows: list[list[str]],
) -> str:
    return "\n".join(
        "\t".join(row)
        for row in rows
        if any(cell for cell in row)
    )


def _docx_table_to_text(
    table: DocxTable,
) -> str:
    rows: list[list[str]] = []

    for row in table.rows:
        cells = [
            normalize_text(cell.text).replace(
                "\n",
                " / ",
            )
            for cell in row.cells
        ]

        rows.append(cells)

    return _table_rows_to_text(rows)


class DOCXParser:
    extensions = frozenset({".docx"})

    def parse(self, path: Path) -> ParseResult:
        try:
            document = DocxDocument(str(path))
        except Exception as exc:
            raise DocumentParsingError(
                "The DOCX file could not be read"
            ) from exc

        sections: list[ParsedUnit] = []

        current_title = "Document"
        current_heading_level: int | None = None
        current_blocks: list[str] = []
        current_block_count = 0

        def flush_section() -> None:
            nonlocal current_blocks
            nonlocal current_block_count

            if not current_blocks and sections:
                return

            content_parts: list[str] = []

            if current_title != "Document":
                content_parts.append(current_title)

            content_parts.extend(current_blocks)

            sections.append(
                create_parsed_unit(
                    unit_index=len(sections) + 1,
                    unit_type="section",
                    source_label=current_title,
                    content="\n\n".join(content_parts),
                    metadata={
                        "heading": current_title,
                        "heading_level": current_heading_level,
                        "block_count": current_block_count,
                    },
                )
            )

            current_blocks = []
            current_block_count = 0

        try:
            blocks = document.iter_inner_content()
        except AttributeError:
            blocks = [
                *document.paragraphs,
                *document.tables,
            ]

        for block in blocks:
            if isinstance(block, Paragraph):
                paragraph_text = normalize_text(
                    block.text
                )

                if not paragraph_text:
                    continue

                style_name = ""

                if block.style is not None:
                    style_name = (
                        block.style.name or ""
                    ).strip()

                if style_name.lower().startswith(
                    "heading"
                ):
                    if current_blocks:
                        flush_section()

                    current_title = paragraph_text

                    level_text = style_name.lower().replace(
                        "heading",
                        "",
                    ).strip()

                    current_heading_level = (
                        int(level_text)
                        if level_text.isdigit()
                        else None
                    )
                else:
                    current_blocks.append(
                        paragraph_text
                    )
                    current_block_count += 1

            elif isinstance(block, DocxTable):
                table_text = _docx_table_to_text(
                    block
                )

                if table_text:
                    current_blocks.append(
                        table_text
                    )
                    current_block_count += 1

        if current_blocks or not sections:
            flush_section()

        return ParseResult(
            parser_name="python-docx",
            file_extension=".docx",
            units=tuple(sections),
            page_count=None,
            word_count=sum(
                unit.word_count for unit in sections
            ),
            quality_score=calculate_quality_score(
                sections
            ),
            requires_ocr=False,
        )


def _pptx_table_to_text(table: Any) -> str:
    rows: list[list[str]] = []

    for row in table.rows:
        row_values = [
            normalize_text(cell.text).replace(
                "\n",
                " / ",
            )
            for cell in row.cells
        ]

        rows.append(row_values)

    return _table_rows_to_text(rows)


def _extract_pptx_shape_text(
    shape: Any,
) -> list[str]:
    extracted_blocks: list[str] = []

    if getattr(shape, "has_text_frame", False):
        paragraphs = [
            normalize_text(paragraph.text)
            for paragraph in shape.text_frame.paragraphs
        ]

        text = "\n".join(
            paragraph
            for paragraph in paragraphs
            if paragraph
        )

        if text:
            extracted_blocks.append(text)

    if getattr(shape, "has_table", False):
        table_text = _pptx_table_to_text(
            shape.table
        )

        if table_text:
            extracted_blocks.append(table_text)

    child_shapes = getattr(shape, "shapes", None)

    if child_shapes is not None:
        for child_shape in child_shapes:
            extracted_blocks.extend(
                _extract_pptx_shape_text(
                    child_shape
                )
            )

    return extracted_blocks


class PPTXParser:
    extensions = frozenset({".pptx"})

    def parse(self, path: Path) -> ParseResult:
        try:
            presentation = Presentation(str(path))
        except Exception as exc:
            raise DocumentParsingError(
                "The PPTX file could not be read"
            ) from exc

        units: list[ParsedUnit] = []
        slide_count = len(presentation.slides)

        for slide_number, slide in enumerate(
            presentation.slides,
            start=1,
        ):
            blocks: list[str] = []

            for shape in slide.shapes:
                blocks.extend(
                    _extract_pptx_shape_text(shape)
                )

            units.append(
                create_parsed_unit(
                    unit_index=slide_number,
                    unit_type="slide",
                    source_label=f"Slide {slide_number}",
                    content="\n\n".join(blocks),
                    metadata={
                        "slide_number": slide_number,
                        "shape_count": len(slide.shapes),
                    },
                )
            )

        text_slide_count = sum(
            1 for unit in units if unit.char_count >= 20
        )

        text_coverage = (
            text_slide_count / slide_count
            if slide_count
            else 0.0
        )

        requires_ocr = bool(
            slide_count and text_coverage < 0.4
        )

        return ParseResult(
            parser_name="python-pptx",
            file_extension=".pptx",
            units=tuple(units),
            page_count=slide_count,
            word_count=sum(
                unit.word_count for unit in units
            ),
            quality_score=calculate_quality_score(
                units,
                expected_unit_count=slide_count,
            ),
            requires_ocr=requires_ocr,
        )


def _spreadsheet_value_to_text(
    value: Any,
) -> str:
    if value is None:
        return ""

    if isinstance(
        value,
        (datetime, date, time),
    ):
        return value.isoformat()

    return normalize_text(str(value)).replace(
        "\n",
        " / ",
    )


class XLSXParser:
    extensions = frozenset({".xlsx"})

    def parse(self, path: Path) -> ParseResult:
        try:
            workbook = load_workbook(
                filename=str(path),
                read_only=True,
                data_only=True,
                keep_links=False,
            )
        except Exception as exc:
            raise DocumentParsingError(
                "The XLSX file could not be read"
            ) from exc

        units: list[ParsedUnit] = []

        try:
            for sheet_index, worksheet in enumerate(
                workbook.worksheets,
                start=1,
            ):
                extracted_rows: list[list[str]] = []
                truncated = False

                declared_max_row = worksheet.max_row or 0
                declared_max_column = (
                    worksheet.max_column or 0
                )

                for row_number, row in enumerate(
                    worksheet.iter_rows(
                        min_row=1,
                        max_row=MAX_SPREADSHEET_ROWS,
                        min_col=1,
                        max_col=min(
                            max(declared_max_column, 1),
                            MAX_SPREADSHEET_COLUMNS,
                        ),
                        values_only=True,
                    ),
                    start=1,
                ):
                    values = [
                        _spreadsheet_value_to_text(
                            value
                        )
                        for value in row
                    ]

                    while values and not values[-1]:
                        values.pop()

                    if any(values):
                        extracted_rows.append(values)

                    if row_number >= MAX_SPREADSHEET_ROWS:
                        truncated = (
                            declared_max_row
                            > MAX_SPREADSHEET_ROWS
                        )
                        break

                if (
                    declared_max_column
                    > MAX_SPREADSHEET_COLUMNS
                ):
                    truncated = True

                units.append(
                    create_parsed_unit(
                        unit_index=sheet_index,
                        unit_type="sheet",
                        source_label=worksheet.title,
                        content=_table_rows_to_text(
                            extracted_rows
                        ),
                        metadata={
                            "sheet_name": worksheet.title,
                            "declared_rows": declared_max_row,
                            "declared_columns": (
                                declared_max_column
                            ),
                            "extracted_rows": len(
                                extracted_rows
                            ),
                            "truncated": truncated,
                        },
                    )
                )

        finally:
            workbook.close()

        return ParseResult(
            parser_name="openpyxl",
            file_extension=".xlsx",
            units=tuple(units),
            page_count=None,
            word_count=sum(
                unit.word_count for unit in units
            ),
            quality_score=calculate_quality_score(
                units
            ),
            requires_ocr=False,
        )